"""
RAG Engine — Extraction, Chunking, Embeddings, Stockage vectoriel Supabase.
"""

import csv
import io
import logging
import os
import re
import time
from datetime import datetime, timezone

import tiktoken
from openai import OpenAI, RateLimitError
from supabase import Client, create_client

from utils.text_cleaner import clean_text

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".txt", ".csv",
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp",
}

# ── OCR (optionnel) ──────────────────────────────────────────────────────────

_tesseract_available = False
try:
    import pytesseract
    from PIL import Image as PILImage
    for _tess_path in [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]:
        if os.path.exists(_tess_path):
            pytesseract.pytesseract.tesseract_cmd = _tess_path
            break
    pytesseract.get_tesseract_version()
    _tesseract_available = True
    logger.info("Tesseract OCR disponible.")
except Exception:
    logger.warning("Tesseract OCR non disponible — images et PDFs scannés ignorés.")


def _ocr_image(image) -> str:
    if not _tesseract_available:
        return ""
    try:
        return pytesseract.image_to_string(image, lang="fra+eng")
    except Exception as exc:
        logger.warning("Échec OCR : %s", exc)
        return ""


# ── Extraction ───────────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> list[dict]:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        raw = page.extract_text() or ""
        cleaned = clean_text(raw)
        if len(cleaned) < 50 and _tesseract_available:
            for img_obj in page.images:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(io.BytesIO(img_obj.data))
                    cleaned += " " + clean_text(_ocr_image(img))
                except Exception:
                    pass
        if cleaned.strip():
            pages.append({"page": i + 1, "text": cleaned.strip()})
    return pages


def _extract_docx(file_bytes: bytes) -> list[dict]:
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))
    GROUP_SIZE = 50
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    for i in range(0, len(paragraphs), GROUP_SIZE):
        block = " ".join(paragraphs[i: i + GROUP_SIZE])
        if block:
            pages.append({"page": i // GROUP_SIZE + 1, "text": block})
    return pages


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        text = " ".join(parts)
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages


def _extract_txt(file_bytes: bytes) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace")
    BLOCK = 3000
    pages = []
    for i in range(0, len(text), BLOCK):
        block = text[i: i + BLOCK].strip()
        if block:
            pages.append({"page": i // BLOCK + 1, "text": block})
    return pages


def _extract_csv(file_bytes: bytes) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(text))
    rows_text = []
    try:
        for row in reader:
            rows_text.append(", ".join(f"{k}: {v}" for k, v in row.items() if v))
    except Exception:
        rows_text = [text]
    GROUP = 100
    pages = []
    for i in range(0, len(rows_text), GROUP):
        block = "\n".join(rows_text[i: i + GROUP])
        if block:
            pages.append({"page": i // GROUP + 1, "text": block})
    return pages


def _extract_image(file_bytes: bytes) -> list[dict]:
    if not _tesseract_available:
        return []
    from PIL import Image as PILImage
    img = PILImage.open(io.BytesIO(file_bytes))
    text = _ocr_image(img).strip()
    return [{"page": 1, "text": text}] if text else []


def is_supported(filename: str) -> bool:
    ext = os.path.splitext(filename.lower())[1]
    return ext in SUPPORTED_EXTENSIONS


def extract_pages(file_bytes: bytes, filename: str) -> list[dict]:
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes)
    elif ext == ".txt":
        return _extract_txt(file_bytes)
    elif ext == ".csv":
        return _extract_csv(file_bytes)
    elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}:
        return _extract_image(file_bytes)
    else:
        raise ValueError(f"Extension non supportée : {ext}")


# ── Chunking ─────────────────────────────────────────────────────────────────

CHUNK_SIZE = 700
OVERLAP = 75
MAX_CHUNK_SIZE = 800


def _split_into_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+", text)
    return [p.strip() for p in parts if p.strip()]


def chunk_pages(pages: list[dict], filename: str, chunk_size: int = CHUNK_SIZE, overlap: int = OVERLAP) -> list[dict]:
    enc = tiktoken.get_encoding("cl100k_base")
    ingested_at = datetime.now(timezone.utc).isoformat()

    pairs: list[tuple[str, int]] = []
    for page_data in pages:
        for sentence in _split_into_sentences(page_data["text"]):
            pairs.append((sentence, page_data["page"]))

    if not pairs:
        return []

    chunks: list[dict] = []
    start = 0

    while start < len(pairs):
        current_sentences: list[str] = []
        current_token_count = 0
        first_page = pairs[start][1]
        end = start

        while end < len(pairs):
            sentence, _ = pairs[end]
            token_count = len(enc.encode(sentence))
            if current_token_count + token_count > MAX_CHUNK_SIZE and current_sentences:
                break
            current_sentences.append(sentence)
            current_token_count += token_count
            end += 1

        if not current_sentences:
            current_sentences = [pairs[start][0]]
            end = start + 1

        chunks.append({
            "content": " ".join(current_sentences),
            "metadata": {
                "filename": filename,
                "page": first_page,
                "chunk_index": len(chunks),
                "ingested_at": ingested_at,
            },
        })

        overlap_token_count = 0
        overlap_start = end
        for i in range(end - 1, start, -1):
            overlap_token_count += len(enc.encode(pairs[i][0]))
            if overlap_token_count >= overlap:
                overlap_start = i
                break

        start = max(start + 1, overlap_start)

    return chunks


# ── Embeddings ───────────────────────────────────────────────────────────────

EMBEDDING_MODEL = "text-embedding-3-small"
BATCH_SIZE = 100
MAX_RETRIES = 3
RETRY_DELAY = 2.0

_openai_client: OpenAI | None = None


def _get_openai_client() -> OpenAI:
    global _openai_client
    if _openai_client is None:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY manquante.")
        _openai_client = OpenAI(api_key=api_key)
    return _openai_client


def generate_embeddings(texts: list[str]) -> list[list[float]]:
    client = _get_openai_client()
    all_embeddings: list[list[float]] = []

    for batch_start in range(0, len(texts), BATCH_SIZE):
        batch = texts[batch_start: batch_start + BATCH_SIZE]
        for attempt in range(MAX_RETRIES):
            try:
                response = client.embeddings.create(model=EMBEDDING_MODEL, input=batch)
                sorted_data = sorted(response.data, key=lambda x: x.index)
                all_embeddings.extend([item.embedding for item in sorted_data])
                break
            except RateLimitError:
                if attempt < MAX_RETRIES - 1:
                    wait = RETRY_DELAY * (2 ** attempt)
                    logger.warning("Rate limit, retry dans %.1fs...", wait)
                    time.sleep(wait)
                else:
                    raise

    return all_embeddings


def generate_embedding(text: str) -> list[float]:
    return generate_embeddings([text])[0]


# ── Stockage Supabase ─────────────────────────────────────────────────────────

SIMILARITY_THRESHOLD = 0.3
TOP_K = 5

_supabase: Client | None = None


def _get_supabase_client() -> Client:
    global _supabase
    if _supabase is None:
        url = os.getenv("SUPABASE_URL")
        key = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
        if not url or not key:
            raise ValueError("SUPABASE_URL ou SUPABASE_SERVICE_ROLE_KEY manquante.")
        _supabase = create_client(url, key)
    return _supabase


def insert_chunks(chunks: list[dict], embeddings: list[list[float]]) -> int:
    client = _get_supabase_client()
    records = [
        {"content": chunk["content"], "embedding": embedding, "metadata": chunk["metadata"]}
        for chunk, embedding in zip(chunks, embeddings)
    ]
    client.table("documents").insert(records).execute()
    logger.info("Inséré %d chunks dans Supabase.", len(records))
    return len(records)


def search_similar_chunks(
    query_embedding: list[float],
    top_k: int = TOP_K,
    threshold: float = SIMILARITY_THRESHOLD,
) -> list[dict]:
    client = _get_supabase_client()
    result = client.rpc(
        "match_documents",
        {"query_embedding": query_embedding, "match_count": top_k, "similarity_threshold": threshold},
    ).execute()
    data = result.data or []
    logger.info("Recherche vectorielle : %d chunk(s) (seuil=%.2f).", len(data), threshold)
    return data


def list_documents() -> list[dict]:
    client = _get_supabase_client()
    result = client.table("documents").select("metadata").execute()
    files: dict[str, dict] = {}
    for row in result.data or []:
        meta = row.get("metadata") or {}
        filename = meta.get("filename", "inconnu")
        if filename not in files:
            files[filename] = {
                "name": filename,
                "chunks": 0,
                "theme": meta.get("theme", "entreprise"),
                "subfolder": meta.get("subfolder"),
            }
        files[filename]["chunks"] += 1
    return sorted(files.values(), key=lambda x: x["name"])


def delete_by_filename(filename: str) -> int:
    client = _get_supabase_client()
    result = (
        client.table("documents")
        .select("id", count="exact")
        .filter("metadata->>filename", "eq", filename)
        .execute()
    )
    count = result.count or 0
    if count > 0:
        client.table("documents").delete().filter("metadata->>filename", "eq", filename).execute()
        logger.info("Supprimé %d chunk(s) pour '%s'.", count, filename)
    return count


def update_file_metadata(filename: str, theme: str | None = None, subfolder: str | None = None) -> int:
    client = _get_supabase_client()
    result = (
        client.table("documents")
        .select("id, metadata")
        .filter("metadata->>filename", "eq", filename)
        .execute()
    )
    count = 0
    for row in result.data or []:
        new_meta = dict(row["metadata"])
        if theme is not None:
            new_meta["theme"] = theme
        if subfolder is not None:
            new_meta["subfolder"] = subfolder if subfolder else None
        client.table("documents").update({"metadata": new_meta}).eq("id", row["id"]).execute()
        count += 1
    logger.info("Metadata mis à jour pour %d chunk(s) de '%s'.", count, filename)
    return count


def delete_all_documents() -> int:
    client = _get_supabase_client()
    result = client.table("documents").select("id", count="exact").execute()
    count = result.count or 0
    if count > 0:
        client.table("documents").delete().neq("id", "00000000-0000-0000-0000-000000000000").execute()
        logger.info("Supprimé %d chunk(s).", count)
    return count
