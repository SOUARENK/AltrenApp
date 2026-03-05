"""
RAG Engine — Moteur de Retrieval-Augmented Generation.

Consolide toutes les opérations RAG :
- Extraction de texte multi-format (PDF, DOCX, PPTX, TXT, CSV, images)
- Chunking token-aware avec overlap
- Génération d'embeddings (OpenAI text-embedding-3-small)
- Stockage et recherche vectorielle (Supabase pgvector)
- Gestion des documents indexés
"""

# ===========================================================================
# SECTION 1 — Extraction de texte (depuis document_service.py)
# ===========================================================================

import csv
import io
import logging
import os
import re
from datetime import datetime, timezone

import tiktoken
from utils.text_cleaner import clean_text

logger = logging.getLogger(__name__)

# Extensions supportées
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".txt", ".csv",
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp",
}

# ---------------------------------------------------------------------------
# OCR — Tesseract (optionnel)
# ---------------------------------------------------------------------------
_tesseract_available = False
try:
    import pytesseract
    from PIL import Image as PILImage

    for _tess_path in [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]:
        if os.path.exists(_tess_path):
            pytesseract.pytesseract.tesseract_cmd = _tess_path
            break

    pytesseract.get_tesseract_version()
    _tesseract_available = True
    logger.info("Tesseract OCR disponible.")
except Exception:
    logger.warning(
        "Tesseract OCR non disponible. "
        "Installez-le depuis https://github.com/UB-Mannheim/tesseract/wiki "
        "pour lire les images et les PDFs scannés."
    )


def _ocr_image(image) -> str:
    """Lance l'OCR sur une image PIL. Retourne '' si Tesseract indisponible."""
    if not _tesseract_available:
        return ""
    try:
        return pytesseract.image_to_string(image, lang="fra+eng")
    except Exception as exc:
        logger.warning("Échec OCR : %s", exc)
        return ""


def _extract_pdf(file_bytes: bytes) -> list[dict]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        raw = page.extract_text() or ""
        cleaned = clean_text(raw)

        if len(cleaned) < 50 and _tesseract_available:
            for img_obj in page.images:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(io.BytesIO(img_obj.data))
                    ocr_text = _ocr_image(img)
                    cleaned += " " + clean_text(ocr_text)
                except Exception:
                    pass

        if cleaned.strip():
            pages.append({"page": i + 1, "text": cleaned.strip()})

    return pages


def _extract_docx(file_bytes: bytes) -> list[dict]:
    import docx

    doc = docx.Document(io.BytesIO(file_bytes))
    GROUP_SIZE = 50
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    for i in range(0, len(paragraphs), GROUP_SIZE):
        block = " ".join(paragraphs[i: i + GROUP_SIZE])
        if block:
            pages.append({"page": i // GROUP_SIZE + 1, "text": block})
    return pages


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = " ".join(parts)
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages


def _extract_txt(file_bytes: bytes) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace")
    BLOCK = 3000
    pages = []
    for i in range(0, len(text), BLOCK):
        block = text[i: i + BLOCK].strip()
        if block:
            pages.append({"page": i // BLOCK + 1, "text": block})
    return pages


def _extract_csv(file_bytes: bytes) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(text))
    rows_text = []
    try:
        for row in reader:
            rows_text.append(", ".join(f"{k}: {v}" for k, v in row.items() if v))
    except Exception:
        rows_text = [text]

    GROUP = 100
    pages = []
    for i in range(0, len(rows_text), GROUP):
        block = "\n".join(rows_text[i: i + GROUP])
        if block:
            pages.append({"page": i // GROUP + 1, "text": block})
    return pages


def _extract_image(file_bytes: bytes) -> list[dict]:
    if not _tesseract_available:
        logger.warning("Image reçue mais Tesseract non disponible — ignorée.")
        return []
    from PIL import Image as PILImage

    img = PILImage.open(io.BytesIO(file_bytes))
    text = _ocr_image(img).strip()
    if text:
        return [{"page": 1, "text": text}]
    return []


def is_supported(filename: str) -> bool:
    """Vérifie si l'extension du fichier est supportée."""
    ext = os.path.splitext(filename.lower())[1]
    return ext in SUPPORTED_EXTENSIONS


def extract_pages(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Extrait le texte d'un fichier en une liste de {page, text}.

    Délègue à l'extracteur approprié selon l'extension.
    """
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes)
    elif ext == ".txt":
        return _extract_txt(file_bytes)
    elif ext == ".csv":
        return _extract_csv(file_bytes)
    elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}:
        return _extract_image(file_bytes)
    else:
        raise ValueError(f"Extension non supportée : {ext}")


# ===========================================================================
# SECTION 2 — Chunking (depuis pdf_service.py)
# ===========================================================================

CHUNK_SIZE = 700
OVERLAP = 75
MAX_CHUNK_SIZE = 800


def _split_into_sentences(text: str) -> list[str]:
    """Découpe un texte en phrases en respectant les fins de phrase."""
    parts = re.split(r"(?<=[.!?])\s+", text)
    return [p.strip() for p in parts if p.strip()]


def chunk_pages(
    pages: list[dict],
    filename: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = OVERLAP,
) -> list[dict]:
    """
    Découpe les pages en chunks de chunk_size tokens avec overlap tokens de chevauchement.

    Ne coupe pas en milieu de phrase. Conserve les métadonnées (fichier, page, index, date).
    """
    enc = tiktoken.get_encoding("cl100k_base")
    ingested_at = datetime.now(timezone.utc).isoformat()

    pairs: list[tuple[str, int]] = []
    for page_data in pages:
        for sentence in _split_into_sentences(page_data["text"]):
            pairs.append((sentence, page_data["page"]))

    if not pairs:
        return []

    chunks: list[dict] = []
    start = 0

    while start < len(pairs):
        current_sentences: list[str] = []
        current_token_count = 0
        first_page = pairs[start][1]
        end = start

        while end < len(pairs):
            sentence, _ = pairs[end]
            token_count = len(enc.encode(sentence))

            if current_token_count + token_count > MAX_CHUNK_SIZE and current_sentences:
                break

            current_sentences.append(sentence)
            current_token_count += token_count
            end += 1

        if not current_sentences:
            current_sentences = [pairs[start][0]]
            end = start + 1

        chunk_content = " ".join(current_sentences)
        chunks.append(
            {
                "content": chunk_content,
                "metadata": {
                    "filename": filename,
                    "page": first_page,
                    "chunk_index": len(chunks),
                    "ingested_at": ingested_at,
                },
            }
        )

        overlap_token_count = 0
        overlap_start = end
        for i in range(end - 1, start, -1):
            overlap_token_count += len(enc.encode(pairs[i][0]))
            if overlap_token_count >= overlap:
                overlap_start = i
                break

        next_start = max(start + 1, overlap_start)
        start = next_start

    return chunks


# ===========================================================================
# SECTION 3 — Embeddings (depuis embedding_service.py)
# ===========================================================================

import time

from openai import OpenAI, RateLimitError

EMBEDDING_MODEL = "text-embedding-3-small"
BATCH_SIZE = 100
MAX_RETRIES = 3
RETRY_DELAY = 2.0

_openai_client: OpenAI | None = None


def _get_openai_client() -> OpenAI:
    global _openai_client
    if _openai_client is None:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY manquante dans les variables d'environnement.")
        _openai_client = OpenAI(api_key=api_key)
    return _openai_client


def generate_embeddings(texts: list[str]) -> list[list[float]]:
    """
    Génère les embeddings pour une liste de textes (batch de BATCH_SIZE).

    Retente automatiquement en cas de RateLimitError avec backoff exponentiel.
    Retourne une liste de vecteurs de dimension 1536.
    """
    client = _get_openai_client()
    all_embeddings: list[list[float]] = []

    for batch_start in range(0, len(texts), BATCH_SIZE):
        batch = texts[batch_start: batch_start + BATCH_SIZE]

        for attempt in range(MAX_RETRIES):
            try:
                response = client.embeddings.create(model=EMBEDDING_MODEL, input=batch)
                sorted_data = sorted(response.data, key=lambda x: x.index)
                all_embeddings.extend([item.embedding for item in sorted_data])
                break
            except RateLimitError:
                if attempt < MAX_RETRIES - 1:
                    wait = RETRY_DELAY * (2 ** attempt)
                    logger.warning("Rate limit atteint, nouvelle tentative dans %.1fs...", wait)
                    time.sleep(wait)
                else:
                    logger.error("Rate limit : toutes les tentatives ont échoué.")
                    raise

    return all_embeddings


def generate_embedding(text: str) -> list[float]:
    """Génère l'embedding d'un seul texte."""
    return generate_embeddings([text])[0]


# ===========================================================================
# SECTION 4 — Stockage vectoriel (depuis vector_service.py)
# ===========================================================================

from supabase import Client, create_client

SIMILARITY_THRESHOLD = 0.3
TOP_K = 5

_supabase: Client | None = None


def _get_supabase_client() -> Client:
    global _supabase
    if _supabase is None:
        url = os.getenv("SUPABASE_URL")
        key = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
        if not url or not key:
            raise ValueError(
                "SUPABASE_URL ou SUPABASE_SERVICE_ROLE_KEY manquante dans les variables d'environnement."
            )
        _supabase = create_client(url, key)
    return _supabase


def insert_chunks(chunks: list[dict], embeddings: list[list[float]]) -> int:
    """
    Insère les chunks et leurs embeddings dans la table documents.

    Retourne le nombre de records insérés.
    """
    client = _get_supabase_client()
    records = [
        {"content": chunk["content"], "embedding": embedding, "metadata": chunk["metadata"]}
        for chunk, embedding in zip(chunks, embeddings)
    ]
    client.table("documents").insert(records).execute()
    logger.info("Inséré %d chunks dans Supabase.", len(records))
    return len(records)


def search_similar_chunks(
    query_embedding: list[float],
    top_k: int = TOP_K,
    threshold: float = SIMILARITY_THRESHOLD,
) -> list[dict]:
    """
    Recherche les top_k chunks les plus similaires à l'embedding de la question.

    Utilise la fonction SQL match_documents (cosinus via pgvector).
    """
    client = _get_supabase_client()
    result = client.rpc(
        "match_documents",
        {"query_embedding": query_embedding, "match_count": top_k, "similarity_threshold": threshold},
    ).execute()
    data = result.data or []
    logger.info("Recherche vectorielle : %d chunk(s) trouvé(s) (seuil=%.2f).", len(data), threshold)
    return data


def list_documents() -> list[dict]:
    """Retourne la liste des fichiers indexés avec thème, sous-dossier et nombre de chunks."""
    client = _get_supabase_client()
    result = client.table("documents").select("metadata").execute()
    files: dict[str, dict] = {}
    for row in result.data or []:
        meta = row.get("metadata") or {}
        filename = meta.get("filename", "inconnu")
        if filename not in files:
            files[filename] = {
                "name": filename,
                "chunks": 0,
                "theme": meta.get("theme", "entreprise"),
                "subfolder": meta.get("subfolder"),
            }
        files[filename]["chunks"] += 1
    return sorted(files.values(), key=lambda x: x["name"])


def delete_by_filename(filename: str) -> int:
    """Supprime tous les chunks d'un fichier spécifique. Retourne le nombre supprimé."""
    client = _get_supabase_client()
    result = (
        client.table("documents")
        .select("id", count="exact")
        .filter("metadata->>filename", "eq", filename)
        .execute()
    )
    count = result.count or 0
    if count > 0:
        client.table("documents").delete().filter("metadata->>filename", "eq", filename).execute()
        logger.info("Supprimé %d chunk(s) pour '%s'.", count, filename)
    return count


def update_file_metadata(filename: str, theme: str | None = None, subfolder: str | None = None) -> int:
    """Met à jour le thème et/ou le sous-dossier de tous les chunks d'un fichier."""
    client = _get_supabase_client()
    result = (
        client.table("documents")
        .select("id, metadata")
        .filter("metadata->>filename", "eq", filename)
        .execute()
    )
    count = 0
    for row in result.data or []:
        new_meta = dict(row["metadata"])
        if theme is not None:
            new_meta["theme"] = theme
        if subfolder is not None:
            new_meta["subfolder"] = subfolder if subfolder else None
        client.table("documents").update({"metadata": new_meta}).eq("id", row["id"]).execute()
        count += 1
    logger.info("Metadata mis à jour pour %d chunk(s) de '%s'.", count, filename)
    return count


def delete_all_documents() -> int:
    """Supprime tous les documents de la table Supabase. Retourne le nombre supprimé."""
    client = _get_supabase_client()
    result = client.table("documents").select("id", count="exact").execute()
    count = result.count or 0
    if count > 0:
        client.table("documents").delete().neq("id", -1).execute()
        logger.info("Supprimé %d chunk(s) de Supabase.", count)
    return count
